# --- Streamlit Configuration MUST BE FIRST ---
import streamlit as st
st.set_page_config(
    page_title="Deepseek Chatbot",
    page_icon="🤖",
    layout="centered",
    initial_sidebar_state="expanded"
)

# --- Standard Library Imports ---
import requests
import json
import os
import logging
import numpy as np

# --- File Processing Imports ---
from PyPDF2 import PdfReader
import pandas as pd
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from pdf2image import convert_from_path

# --- ML/NLP Imports ---
import faiss
import nltk
from sentence_transformers import SentenceTransformer

# --- Initial Configurations ---
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
nltk.download('punkt', quiet=True)

# --- Platform-Specific Configuration ---
if os.name == 'nt':  # Windows only configuration
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
else:  # Linux/Streamlit Cloud configuration
    pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'

# --- Core Functions ---
def categorize_file(file_path):
    """Cross-platform file type detection"""
    extension = os.path.splitext(file_path)[1].lower()
    return {
        '.jpg': 'image', '.jpeg': 'image', '.png': 'image',
        '.gif': 'image', '.bmp': 'image', '.pdf': 'pdf',
        '.csv': 'csv', '.xlsx': 'xlsx', '.pptx': 'pptx',
        '.docx': 'docx', '.txt': 'text', '.md': 'text'
    }.get(extension, 'unknown')

def process_office_file(file_path):
    """Generic office file processing for Linux/Cloud"""
    file_type = categorize_file(file_path)
    content = []
    
    if file_type == 'pptx':
        prs = Presentation(file_path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes])
    
    if file_type == 'docx':
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    
    return "Unsupported file format"

# ... (rest of processing functions remain the same)

# --- Streamlit UI Components ---
st.title("Deepseek Chatbot 💬")
st.caption("Powered by Ollama - Upload documents for contextual understanding")

# Session state initialization
if "messages" not in st.session_state:
    st.session_state.messages = []

if "uploaded_files" not in st.session_state:
    st.session_state.uploaded_files = []

# ... (rest of UI code remains the same)
