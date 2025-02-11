import streamlit as st
import requests
import json
import os
import logging
import numpy as np
from PyPDF2 import PdfReader
import pandas as pd
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from pdf2image import convert_from_path
datetime = __import__('datetime')

# --- Streamlit Configuration ---
st.set_page_config(
    page_title="Deepseek Chatbot",
    page_icon="🤖",
    layout="centered",
    initial_sidebar_state="expanded"
)

# Initialize session state variables
if "messages" not in st.session_state:
    st.session_state.messages = []

if "uploaded_files" not in st.session_state:
    st.session_state.uploaded_files = []

def categorize_file(file_path):
    extension = os.path.splitext(file_path)[1].lower()
    return {
        '.jpg': 'image', '.jpeg': 'image', '.png': 'image',
        '.pdf': 'pdf', '.csv': 'csv', '.xlsx': 'xlsx',
        '.pptx': 'pptx', '.docx': 'docx', '.txt': 'text'
    }.get(extension, 'unknown')

def process_pdf(file_path):
    text = ""
    try:
        reader = PdfReader(file_path)
        text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        if not text.strip():  # Fallback to OCR
            images = convert_from_path(file_path)
            text = " ".join([pytesseract.image_to_string(img) for img in images])
    except Exception as e:
        logging.error(f"PDF processing error: {e}")
    return text

def process_file(file_path):
    file_type = categorize_file(file_path)
    try:
        if file_type == 'image':
            return pytesseract.image_to_string(Image.open(file_path))
        elif file_type == 'pdf':
            return process_pdf(file_path)
        elif file_type == 'csv':
            return pd.read_csv(file_path).to_string()
        elif file_type == 'xlsx':
            return pd.read_excel(file_path).to_string()
        elif file_type == 'pptx':
            prs = Presentation(file_path)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame])
        elif file_type == 'docx':
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        elif file_type == 'text':
            with open(file_path, 'r') as f:
                return f.read()
    except Exception as e:
        logging.error(f"Error processing {file_path}: {e}")
        return None

def send_request(prompt, file_path=None):
    api_url = "http://localhost:11434/api/generate"
    if file_path:
        content = process_file(file_path)
        if content:
            prompt += f"\n\nFile Content:\n{content[:3000]}"
    try:
        response = requests.post(api_url, json={"model": "deepseek-r1:1.5b", "prompt": prompt, "stream": False})
        return response.json().get('response', 'No response generated')
    except Exception as e:
        logging.error(f"API Error: {e}")
        return "Failed to get response from API"

def save_chat():
    if not os.path.exists("saved_chats"):
        os.makedirs("saved_chats")
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    chat_path = f"saved_chats/chat_{timestamp}.json"
    with open(chat_path, "w") as f:
        json.dump(st.session_state.messages, f, indent=4)

def load_past_chats():
    if not os.path.exists("saved_chats"):
        return []
    return sorted(os.listdir("saved_chats"), reverse=True)

def load_chat(chat_file):
    with open(f"saved_chats/{chat_file}", "r") as f:
        st.session_state.messages = json.load(f)

# --- Streamlit UI ---
st.title("Deepseek Chatbot 💬")
st.caption("Powered by Ollama - Upload documents for contextual understanding")

# Sidebar - Past chats
with st.sidebar:
    st.header("Past Chats")
    past_chats = load_past_chats()
    if past_chats:
        selected_chat = st.selectbox("Select a past chat", past_chats, index=0)
        if st.button("Load Chat"):
            load_chat(selected_chat)
    if st.button("End Chat & Save"):
        save_chat()
        st.session_state.messages = []
        st.success("Chat saved successfully!")

# Sidebar - File upload
with st.sidebar:
    st.header("Document Upload")
    uploaded_files = st.file_uploader("Choose files", type=["pdf", "png", "jpg", "jpeg", "csv", "xlsx", "docx", "pptx", "txt"], accept_multiple_files=True)
    if uploaded_files:
        st.session_state.uploaded_files = uploaded_files
        st.success(f"{len(uploaded_files)} file(s) ready for analysis!")

# Display chat history
for message in st.session_state.get("messages", []):
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# Chat input
if prompt := st.chat_input("Ask me anything..."):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)
    with st.chat_message("assistant"):
        response_placeholder = st.empty()
        full_response = ""
        if st.session_state.uploaded_files:
            temp_dir = "temp_files"
            os.makedirs(temp_dir, exist_ok=True)
            for file in st.session_state.uploaded_files:
                file_path = os.path.join(temp_dir, file.name)
                with open(file_path, "wb") as f:
                    f.write(file.getvalue())
                try:
                    response = send_request(prompt, file_path)
                    full_response += f"**{file.name} Analysis:**\n{response}\n\n"
                finally:
                    os.remove(file_path)
        else:
            full_response = send_request(prompt)
        response_placeholder.markdown(full_response)
    st.session_state.messages.append({"role": "assistant", "content": full_response})
